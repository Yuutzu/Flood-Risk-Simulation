"""
Fix Jade Valley PPT — focus on flood_animation.py as main system,
flood_risk_simulation.py as supporting HAND model only.
Remove Arena, fix factual errors, make rubric honest.
"""
import os
from pptx import Presentation

SRC = os.path.expanduser("~/Downloads/Jade_Valley_Flood_Risk_Simulation.pptx")
DST = os.path.expanduser("~/Downloads/Jade_Valley_Flood_Risk_Simulation_FIXED2.pptx")

prs = Presentation(SRC)
slides = list(prs.slides)


def replace_all(slide, old, new):
    """Replace text across all shapes in a slide, preserving first run's formatting."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full_text = para.text
            if old in full_text:
                new_text = full_text.replace(old, new)
                if para.runs:
                    first_run = para.runs[0]
                    fn = first_run.font.name
                    fs = first_run.font.size
                    fb = first_run.font.bold
                    try:
                        fc = first_run.font.color.rgb
                    except Exception:
                        fc = None
                    for run in para.runs:
                        run.text = ""
                    para.runs[0].text = new_text
                    para.runs[0].font.name = fn
                    if fs:
                        para.runs[0].font.size = fs
                    para.runs[0].font.bold = fb
                    if fc:
                        para.runs[0].font.color.rgb = fc


# SLIDE 1 — Title
s1 = slides[0]
replace_all(s1,
    "HAND Model  \u2022  PAGASA Return Periods  \u2022  Arena-Based Discrete-Event Simulation",
    "Dynamic Flood Animation  \u2022  HAND Risk Model  \u2022  PAGASA Scenarios  \u2022  Python Geospatial Simulation")

# SLIDE 2 — Problem Definition & Objectives
s2 = slides[1]
replace_all(s2, "EXCELLENT", "")
replace_all(s2, "15 / 15 pts", "")
replace_all(s2,
    "Model water flow and accumulation across terrain during heavy rainfall and storm surge events using Arena Simulation",
    "Simulate real-time flood progression using D8 flow routing and river overflow physics in an interactive animated viewer")
replace_all(s2,
    "Identify high-risk flood zones and areas of dangerous water accumulation within the simulated environment",
    "Identify high-risk flood zones using the HAND model and classify risk across 5 severity levels")
replace_all(s2,
    "Measure simulated flood severity, inundation depth, and spread under different rainfall intensities",
    "Visualize flood depth, river overflow, and area inundation across 6 PAGASA-based rainfall scenarios")
replace_all(s2,
    "Evaluate effectiveness of infrastructure interventions \u2014 drainage canals, retention basins, and elevated barriers",
    "Provide an interactive GUI with preset, manual, and randomized scenario selection for exploratory analysis")
replace_all(s2,
    "Provide data-driven insights and recommendations to support urban planning and flood mitigation decision-making",
    "Generate evacuation alerts (NORMAL \u2192 EVACUATE NOW) and export GIF animations for offline review and planning")

# SLIDE 3 — Model Design
s3 = slides[2]
replace_all(s3, "EXCELLENT", "")
replace_all(s3, "Full Credit", "")
replace_all(s3,
    "Each 30.64m\u00b2 grid cell is modeled as a distinct entity carrying elevation, slope, and HAND value attributes",
    "Each 30.64 m grid cell carries elevation, slope, flow accumulation, saturation, rain water depth, and river water depth")
replace_all(s3,
    "Stream channels derived from top 5% flow accumulation \u2014 the primary resource limiting flood conveyance",
    "Two-layer water system: rain runoff (D8 routing, 12 iterations) + river overflow (BFS dilation with accumulation ramp)")
replace_all(s3,
    "Height Above Nearest Drainage computed per cell using D8 flow routing and ascending elevation propagation",
    "Per-timestep cycle: Rain \u2192 D8 Route \u2192 River Overflow \u2192 Infiltration \u2192 Drainage, repeated each 10-min step")
replace_all(s3,
    "Each cell tagged with risk level (Very High \u2192 Safe) based on HAND threshold for each return period scenario",
    "Real-time alert classification: NORMAL \u2192 STANDBY \u2192 PRE-EVACUATION \u2192 MANDATORY EVACUATION \u2192 EVACUATE NOW")
replace_all(s3,
    "Checks HAND value against return-period thresholds (1.0m, 2.0m, 3.5m, 6.0m) to route cell to correct zone",
    "6 scenario presets (Light Rain to Typhoon Signal 3) + Custom mode with 7 adjustable parameters")
replace_all(s3,
    "Inundated cells exit system; statistics (ha flooded, depth, % area) recorded for each return period",
    "Interactive matplotlib viewer with play/pause, scrubber, speed control + GIF export to Results/animations/")
replace_all(s3, "ENTITY", "GRID")
replace_all(s3, "Terrain Cells", "57\u00d761 DEM Cells")
replace_all(s3, "RESOURCE", "PHYSICS")
replace_all(s3, "Drainage Network", "Two-Layer Model")
replace_all(s3, "PROCESS", "TIMESTEP")
replace_all(s3, "HAND Computation", "Step Cycle")
replace_all(s3, "ASSIGN", "ALERTS")
replace_all(s3, "Risk Classification", "Evacuation Levels")
replace_all(s3, "DECIDE", "SCENARIOS")
replace_all(s3, "Flood Scenario", "Preset / Custom")
replace_all(s3, "DISPOSE", "OUTPUT")
replace_all(s3, "Output / Report", "Viewer + GIF")

# SLIDE 4 — Process Flow
s4 = slides[3]
replace_all(s4, "EXCELLENT", "")
replace_all(s4, "Full Credit", "")
replace_all(s4, "Arena Modules:", "Python Pipeline:")
replace_all(s4, "CREATE", "LOAD DEM")
replace_all(s4, "ASSIGN", "HYDRO")
replace_all(s4, "DECIDE", "SIMULATE")
replace_all(s4, "PROCESS", "ANIMATE")
replace_all(s4, "DISPOSE", "EXPORT")
replace_all(s4, "Load GeoTIFF + DXF", "Load GeoTIFF DEM")
replace_all(s4, "Fill NoData cells", "Fill NoData + smooth")
replace_all(s4, "Smoothing (\u03c3=0.5)", "Build flow grid")
replace_all(s4, "Sink filling", "D8 flow direction")
replace_all(s4, "D8 Flow Direction", "Flow accumulation")
replace_all(s4, "Flow Accumulation", "Stream mask (top 8%)")
replace_all(s4, "Top 5% flow acc.", "River channel BFS")
replace_all(s4, "Network extraction", "Bank elevation calc")
replace_all(s4, "Channel mapping", "Overflow thresholds")
replace_all(s4, "Vertical distance", "Rain + Route + Overflow")
replace_all(s4, "to nearest drainage", "+ Infiltration + Drain")
replace_all(s4, "along flow path", "per 10-min timestep")
replace_all(s4, "4 return periods:", "6 presets + Custom")
replace_all(s4, "5 / 10 / 25 / 100 yr", "Light Rain \u2192 Signal 3")
replace_all(s4, "Threshold mapping", "tkinter GUI launcher")
replace_all(s4, "5-class risk map", "matplotlib viewer")
replace_all(s4, "Statistics output", "Play/Pause/Scrub/GIF")
replace_all(s4, "Report generation", "Evacuation alerts")

# SLIDE 5 — Terrain (from static model)
s5 = slides[4]
replace_all(s5, "EXCELLENT", "")
replace_all(s5, "Full Credit", "")
replace_all(s5, "Simulation Setup & Terrain Analysis",
    "Terrain Data & HAND Risk Model (flood_risk_simulation.py)")

# SLIDE 6 — Flood Scenarios
s6 = slides[5]
replace_all(s6, "EXCELLENT", "")
replace_all(s6, "Full Credit", "")
replace_all(s6,
    "Flood Scenarios \u2014 PAGASA Return Periods",
    "Dynamic Scenarios \u2014 PAGASA-Based Presets (flood_animation.py)")
replace_all(s6,
    "Based on PAGASA standards",
    "Rainfall rates based on PAGASA classification (mm/hr ranges)")

# SLIDE 7 — Risk Zones
s7 = slides[6]
replace_all(s7, "EXCELLENT", "")
replace_all(s7, "Full Credit", "")

# SLIDE 10 — Rubric
s10 = slides[9]
replace_all(s10,
    "Complete use of CREATE, ASSIGN, DECIDE, PROCESS, DISPOSE; logical HAND model flow",
    "Two-layer physics model (D8 + BFS); pipeline: LOAD \u2192 HYDRO \u2192 SIMULATE \u2192 ANIMATE \u2192 EXPORT")
replace_all(s10,
    "All 5 criteria rated EXCELLENT  \u2022  Maximum rubric score achieved across all dimensions",
    "Self-assessment based on rubric criteria  \u2022  Scores reflect honest evaluation of deliverables")

# SLIDE 11 — Thank You
s11 = slides[10]
replace_all(s11,
    "Simulation Method: Arena Discrete-Event Simulation",
    "Simulation: Python Dynamic Flood Animation + HAND Risk Model")

# SAVE
prs.save(DST)
print(f"\nSaved to:\n  {DST}")
print()
print("AUTO-FIXED:")
print("  OK  All Arena references -> Python Dynamic Flood Animation")
print("  OK  All self-grade labels removed (EXCELLENT / Full Credit / 15/15)")
print("  OK  Slide 2 objectives rewritten for flood_animation.py focus")
print("  OK  Slide 3 model design -> two-layer physics + alerts + GUI")
print("  OK  Slide 4 process flow -> animation pipeline steps")
print("  OK  Slide 5 title -> credits flood_risk_simulation.py as model source")
print("  OK  Slide 6 title -> dynamic scenarios from flood_animation.py")
print("  OK  Fixed '30.64m2' -> '30.64 m grid cell'")
print("  OK  Fixed PAGASA claim -> 'PAGASA classification (mm/hr ranges)'")
print("  OK  Rubric -> honest language, animation pipeline reference")
