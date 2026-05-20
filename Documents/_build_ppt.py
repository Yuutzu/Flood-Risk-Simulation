"""Generate Jade_Valley_Flood_Simulation_FINAL.pptx — monochrome formal design.

Run from repo root:
  python "Documents/_build_ppt.py"
"""
from __future__ import annotations

import csv
import json
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# =============================================================================
# DESIGN TOKENS  (monochrome, formal, brutalist)
# =============================================================================

BLACK   = RGBColor(0x00, 0x00, 0x00)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREY_D  = RGBColor(0x33, 0x33, 0x33)
GREY_M  = RGBColor(0x66, 0x66, 0x66)
GREY_L  = RGBColor(0xCC, 0xCC, 0xCC)
GREY_XL = RGBColor(0xEE, 0xEE, 0xEE)

FONT     = "Calibri"
FONT_MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Margins
MARGIN_L = Inches(0.6)
MARGIN_R = Inches(0.6)
MARGIN_T = Inches(0.55)
MARGIN_B = Inches(0.5)

# =============================================================================
# PATHS
# =============================================================================

BASE     = Path(__file__).resolve().parent.parent
DOCS_DIR = BASE / "Documents"
DATA_DIR = BASE / "Results" / "data"
MAPS_DIR = BASE / "Results" / "maps"
OUT_PPT  = DOCS_DIR / "Jade_Valley_Flood_Simulation_FINAL.pptx"


# =============================================================================
# HELPERS
# =============================================================================

def _set_run(run, text, *, size=12, bold=False, mono=False, color=BLACK,
             italic=False):
    run.text = text
    run.font.name = FONT_MONO if mono else FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text="", *,
                size=12, bold=False, mono=False, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, color=BLACK, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    _set_run(p.add_run(), text, size=size, bold=bold, mono=mono,
             color=color, italic=italic)
    return tb


def add_paragraph_run(text_frame, text, *, size=12, bold=False, mono=False,
                      align=PP_ALIGN.LEFT, color=BLACK, italic=False,
                      space_before=None, space_after=None,
                      first_paragraph=False):
    if first_paragraph and len(text_frame.paragraphs) == 1 \
       and not text_frame.paragraphs[0].text:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()
    p.alignment = align
    if space_before is not None: p.space_before = Pt(space_before)
    if space_after  is not None: p.space_after  = Pt(space_after)
    _set_run(p.add_run(), text, size=size, bold=bold, mono=mono,
             color=color, italic=italic)
    return p


def add_rect(slide, left, top, width, height, *, fill=BLACK, line=None,
             line_w=Pt(0)):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.fill.background() if line is None else None
    if line is not None:
        rect.line.color.rgb = line
        rect.line.width = line_w
    else:
        rect.line.width = Pt(0)
    rect.shadow.inherit = False
    return rect


def add_line(slide, x1, y1, x2, y2, *, color=BLACK, width=Pt(1)):
    line = slide.shapes.add_connector(1, x1, y1, x2 - x1, y2 - y1)
    line.line.color.rgb = color
    line.line.width = width
    # Force start at (x1,y1), end at (x2,y2)
    line.left, line.top = x1, y1
    line.width, line.height = (x2 - x1), (y2 - y1)
    return line


def add_slide_header(slide, section_num, section_title, slide_num, total):
    """Brutalist header block: thick black bar + section/title text."""
    # Top black bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.45), fill=BLACK)
    # Section indicator (left)
    add_textbox(slide, Inches(0.4), Inches(0.05), Inches(8), Inches(0.4),
                f"{section_num}   {section_title.upper()}",
                size=14, bold=True, color=WHITE)
    # Slide counter (right)
    add_textbox(slide, Inches(11.0), Inches(0.05), Inches(2.0), Inches(0.4),
                f"{slide_num} / {total}",
                size=11, bold=True, color=WHITE, mono=True,
                align=PP_ALIGN.RIGHT)
    # Bottom rule
    add_rect(slide, Inches(0.6), Inches(7.05), SLIDE_W - Inches(1.2), Emu(9525),
             fill=BLACK)
    # Footer text
    add_textbox(slide, Inches(0.6), Inches(7.10), Inches(8), Inches(0.3),
                "JADE VALLEY FLOOD SIMULATION   |   DAVAO CITY, PHILIPPINES",
                size=8, bold=True, color=GREY_M, mono=True)
    add_textbox(slide, Inches(10.0), Inches(7.10), Inches(2.7), Inches(0.3),
                f"FINAL BUILD 100%",
                size=8, bold=True, color=GREY_M, mono=True,
                align=PP_ALIGN.RIGHT)


def add_slide_title(slide, title):
    """Big slide title under the header bar."""
    add_textbox(slide, Inches(0.6), Inches(0.65), SLIDE_W - Inches(1.2), Inches(0.6),
                title.upper(), size=24, bold=True, color=BLACK)
    # Thin rule under title
    add_rect(slide, Inches(0.6), Inches(1.20), Inches(2.0), Emu(28575), fill=BLACK)


# =============================================================================
# SLIDE FACTORIES
# =============================================================================

def make_blank(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def slide_cover(prs, total_slides):
    slide = make_blank(prs)
    # Big top black band
    add_rect(slide, 0, 0, SLIDE_W, Inches(3.4), fill=BLACK)
    # Bottom black band
    add_rect(slide, 0, Inches(7.0), SLIDE_W, Inches(0.5), fill=BLACK)
    # Vertical bar (left)
    add_rect(slide, Inches(0.6), Inches(0.4), Inches(0.08), Inches(2.6),
             fill=WHITE)

    add_textbox(slide, Inches(0.95), Inches(0.55), Inches(11.5), Inches(0.4),
                "FINAL PROJECT  |  SIMULATION SYSTEM",
                size=12, bold=True, color=WHITE, mono=True)
    add_textbox(slide, Inches(0.95), Inches(1.0), Inches(11.5), Inches(1.2),
                "JADE VALLEY",
                size=66, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.95), Inches(1.95), Inches(11.5), Inches(1.0),
                "FLOOD SIMULATION",
                size=44, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.95), Inches(2.75), Inches(11.5), Inches(0.4),
                "DAVAO CITY, PHILIPPINES",
                size=14, bold=True, color=WHITE, mono=True)

    # Subtitle below the dark band
    add_textbox(slide, Inches(0.6), Inches(3.7), Inches(11.5), Inches(0.5),
                "A TERRAIN-AWARE HYDROLOGICAL SIMULATION SYSTEM",
                size=18, bold=True, color=BLACK)
    add_textbox(slide, Inches(0.6), Inches(4.15), Inches(11.5), Inches(0.5),
                "D8 FLOW ROUTING  |  GREEN-AMPT INFILTRATION  |  BFS RIVER OVERFLOW",
                size=12, bold=True, color=GREY_D, mono=True)

    # Specs panel (table-like)
    panel_top = Inches(4.85)
    panel_h   = Inches(2.0)
    add_rect(slide, Inches(0.6), panel_top, Inches(12.1), Inches(0.04),
             fill=BLACK)
    add_rect(slide, Inches(0.6), panel_top + panel_h, Inches(12.1), Inches(0.04),
             fill=BLACK)

    rows = [
        ("BUILDS",      "25%  |  50%  |  75%  |  100% (master)"),
        ("SCENARIOS",   "6 PAGASA-calibrated storm scenarios"),
        ("CONFIGS",     "4 prevention configurations (baseline + 3 mitigation levels)"),
        ("MATRIX",      "24 simulation runs in the master dataset"),
        ("STUDY AREA",  "326.44 ha  |  57 x 61 grid  |  ~30.6 m / cell"),
        ("DATE",        datetime.now().strftime("%B %Y")),
    ]
    row_h = Inches(0.30)
    y = panel_top + Inches(0.10)
    for label, value in rows:
        add_textbox(slide, Inches(0.85), y, Inches(2.5), row_h,
                    label, size=10, bold=True, color=BLACK, mono=True)
        add_textbox(slide, Inches(3.35), y, Inches(9.0), row_h,
                    value, size=11, color=BLACK)
        y += row_h

    add_textbox(slide, Inches(0.6), Inches(7.10), Inches(8), Inches(0.3),
                "DOCUMENTATION + RESULTS  |  COMPILED FROM SIMULATION SYSTEM OUTPUTS",
                size=8, bold=True, color=WHITE, mono=True)
    add_textbox(slide, Inches(10.0), Inches(7.10), Inches(2.7), Inches(0.3),
                f"01 / {total_slides:02d}",
                size=8, bold=True, color=WHITE, mono=True,
                align=PP_ALIGN.RIGHT)


def slide_toc(prs, sections, slide_num, total):
    slide = make_blank(prs)
    add_slide_header(slide, "00", "CONTENTS", slide_num, total)
    add_slide_title(slide, "Table of Contents")

    # Two-column TOC
    n = len(sections)
    half = (n + 1) // 2
    col_w = Inches(5.8)
    row_h = Inches(0.42)

    def draw_col(items, left, start_idx):
        y = Inches(1.6)
        for i, (num, title) in enumerate(items):
            idx = start_idx + i
            # Number in heavy black
            add_textbox(slide, left, y, Inches(0.7), row_h,
                        f"{num}", size=18, bold=True, color=BLACK, mono=True)
            # Title
            add_textbox(slide, left + Inches(0.85), y, col_w - Inches(0.85), row_h,
                        title, size=14, color=BLACK,
                        anchor=MSO_ANCHOR.MIDDLE)
            # Dot rule
            add_rect(slide, left, y + Inches(0.40), col_w - Inches(0.1),
                     Emu(9525), fill=GREY_L)
            y += row_h + Inches(0.05)

    draw_col(sections[:half], Inches(0.6), 1)
    draw_col(sections[half:], Inches(7.0), half + 1)


def slide_section_divider(prs, section_num, title, blurb,
                          slide_num, total):
    slide = make_blank(prs)
    # Massive section block: full-bleed black left third, content on right
    add_rect(slide, 0, 0, Inches(4.3), SLIDE_H, fill=BLACK)
    # Section number (giant)
    add_textbox(slide, Inches(0.4), Inches(1.0), Inches(3.7), Inches(2.5),
                section_num,
                size=140, bold=True, color=WHITE, mono=True)
    add_textbox(slide, Inches(0.4), Inches(3.6), Inches(3.7), Inches(0.5),
                "SECTION", size=14, bold=True, color=WHITE, mono=True)

    # Right side: title and blurb
    add_textbox(slide, Inches(4.7), Inches(1.5), Inches(8.4), Inches(1.5),
                title.upper(), size=42, bold=True, color=BLACK)
    add_rect(slide, Inches(4.7), Inches(3.05), Inches(2.0), Emu(28575),
             fill=BLACK)
    add_textbox(slide, Inches(4.7), Inches(3.2), Inches(8.4), Inches(3.5),
                blurb, size=14, color=GREY_D)

    add_textbox(slide, Inches(4.7), Inches(7.10), Inches(6), Inches(0.3),
                "JADE VALLEY FLOOD SIMULATION   |   DAVAO CITY",
                size=8, bold=True, color=GREY_M, mono=True)
    add_textbox(slide, Inches(10.0), Inches(7.10), Inches(2.7), Inches(0.3),
                f"{slide_num} / {total}",
                size=8, bold=True, color=GREY_M, mono=True,
                align=PP_ALIGN.RIGHT)


def slide_text_content(prs, section_num, section_title, slide_title,
                       paragraphs, slide_num, total, *,
                       body_size=13, bullet=False, mono_body=False):
    """paragraphs: list of strings OR list of (label, text) tuples."""
    slide = make_blank(prs)
    add_slide_header(slide, section_num, section_title, slide_num, total)
    add_slide_title(slide, slide_title)

    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.45),
                                    SLIDE_W - Inches(1.2),
                                    Inches(5.55))
    tf = body.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)

    first = True
    for item in paragraphs:
        if isinstance(item, tuple) and len(item) == 2:
            label, txt = item
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(6)
            _set_run(p.add_run(), label + "   ", size=body_size, bold=True,
                     color=BLACK, mono=True)
            _set_run(p.add_run(), txt, size=body_size, color=BLACK,
                     mono=mono_body)
        else:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(6)
            text = item
            if bullet and text.strip():
                text = "—  " + text
            _set_run(p.add_run(), text, size=body_size, color=BLACK,
                     mono=mono_body)


def slide_table(prs, section_num, section_title, slide_title,
                headers, rows, slide_num, total, *,
                col_widths=None, header_size=10, body_size=9, mono=True,
                caption=None):
    slide = make_blank(prs)
    add_slide_header(slide, section_num, section_title, slide_num, total)
    add_slide_title(slide, slide_title)

    n_cols = len(headers)
    n_rows = len(rows) + 1
    table_w = SLIDE_W - Inches(1.2)
    if col_widths is None:
        col_widths = [table_w // n_cols] * n_cols

    table_top = Inches(1.45)
    row_h = Inches(0.32) if n_rows <= 14 else Inches(0.26)
    if row_h * n_rows > Inches(5.3):
        row_h = Inches(5.3) // n_rows

    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.6), table_top, table_w,
        row_h * n_rows).table

    for i, w in enumerate(col_widths):
        table_shape.columns[i].width = w
    for r in range(n_rows):
        table_shape.rows[r].height = row_h

    # Header row
    for c, h in enumerate(headers):
        cell = table_shape.cell(0, c)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLACK
        tf = cell.text_frame
        tf.margin_left = Emu(36000); tf.margin_right = Emu(36000)
        tf.margin_top = Emu(18000); tf.margin_bottom = Emu(18000)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        _set_run(p.add_run(), str(h), size=header_size, bold=True,
                 color=WHITE, mono=mono)

    # Body rows
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table_shape.cell(r, c)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else GREY_XL
            tf = cell.text_frame
            tf.margin_left = Emu(36000); tf.margin_right = Emu(36000)
            tf.margin_top = Emu(14000); tf.margin_bottom = Emu(14000)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            _set_run(p.add_run(), str(val), size=body_size, color=BLACK,
                     mono=mono)

    if caption:
        add_textbox(slide, Inches(0.6),
                    table_top + row_h * n_rows + Inches(0.1),
                    SLIDE_W - Inches(1.2), Inches(0.5),
                    caption, size=10, color=GREY_M, italic=True)


def slide_raw_text(prs, section_num, section_title, slide_title,
                   text_path, slide_num, total, *, caption=None,
                   body_size=7.5):
    """Embed a raw text file verbatim as a monospace block (preserves layout)."""
    slide = make_blank(prs)
    add_slide_header(slide, section_num, section_title, slide_num, total)
    add_slide_title(slide, slide_title)

    body_top = Inches(1.45)
    body_h   = Inches(5.4)
    if caption:
        body_h = Inches(5.0)

    # Black border box around the text for "brutal" framing
    add_rect(slide, Inches(0.6), body_top, SLIDE_W - Inches(1.2), body_h,
             fill=WHITE)
    # Thin border
    border_w = Pt(1.25)
    add_rect(slide, Inches(0.6), body_top,
             SLIDE_W - Inches(1.2), Emu(12700), fill=BLACK)
    add_rect(slide, Inches(0.6), body_top + body_h - Emu(12700),
             SLIDE_W - Inches(1.2), Emu(12700), fill=BLACK)
    add_rect(slide, Inches(0.6), body_top,
             Emu(12700), body_h, fill=BLACK)
    add_rect(slide, SLIDE_W - Inches(0.6) - Emu(12700), body_top,
             Emu(12700), body_h, fill=BLACK)

    tb = slide.shapes.add_textbox(
        Inches(0.75), body_top + Inches(0.10),
        SLIDE_W - Inches(1.5), body_h - Inches(0.20))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top  = Emu(0); tf.margin_bottom = Emu(0)

    text = Path(text_path).read_text(encoding="utf-8") if Path(text_path).exists() \
           else f"[File not found: {text_path}]"
    lines = text.splitlines() or [""]

    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0); p.space_after = Pt(0)
        _set_run(p.add_run(), line or " ",
                 size=body_size, mono=True, color=BLACK)

    if caption:
        add_textbox(slide, Inches(0.6), body_top + body_h + Inches(0.1),
                    SLIDE_W - Inches(1.2), Inches(0.4),
                    caption, size=10, color=GREY_M, italic=True)


def slide_image(prs, section_num, section_title, slide_title, image_path,
                slide_num, total, caption=None):
    slide = make_blank(prs)
    add_slide_header(slide, section_num, section_title, slide_num, total)
    add_slide_title(slide, slide_title)
    if image_path and Path(image_path).exists():
        # Centre-fit image into remaining canvas
        max_w = SLIDE_W - Inches(1.2)
        max_h = Inches(5.0)
        pic = slide.shapes.add_picture(str(image_path),
                                       Inches(0.6), Inches(1.45),
                                       width=max_w)
        if pic.height > max_h:
            ratio = max_h / pic.height
            pic.width  = int(pic.width * ratio)
            pic.height = max_h
            pic.left   = (SLIDE_W - pic.width) // 2
        else:
            pic.left   = (SLIDE_W - pic.width) // 2
    else:
        add_textbox(slide, Inches(0.6), Inches(3), SLIDE_W - Inches(1.2),
                    Inches(0.5), f"[Image not found: {image_path}]",
                    size=12, color=GREY_M, italic=True,
                    align=PP_ALIGN.CENTER)
    if caption:
        add_textbox(slide, Inches(0.6), Inches(6.55), SLIDE_W - Inches(1.2),
                    Inches(0.4), caption, size=10, color=GREY_M,
                    italic=True, align=PP_ALIGN.CENTER)


# =============================================================================
# DATA LOADING
# =============================================================================

def load_master_csv():
    rows = []
    path = DATA_DIR / "master_results_table.csv"
    if not path.exists():
        return rows
    with open(path, newline="", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        for r in reader:
            rows.append(r)
    return rows


def load_flood_statistics():
    path = DATA_DIR / "flood_statistics.json"
    if not path.exists():
        return None
    with open(path, encoding="utf-8") as f:
        return json.load(f)


# =============================================================================
# PRESENTATION ASSEMBLY
# =============================================================================

# Section structure (number, short title, divider blurb)
SECTIONS = [
    ("01", "Introduction",
     "Project context, geographic setting, and an overview of the three "
     "progressive simulation builds plus the 100% master orchestrator."),
    ("02", "Problem Definition",
     "Why Jade Valley is flood-prone, the planning gap this system addresses, "
     "and the critical questions a terrain-aware model can answer."),
    ("03", "Objectives",
     "Six concrete objectives the simulation system is designed to meet."),
    ("04", "System Description",
     "Architecture, build hierarchy, core components, and data flow."),
    ("05", "System Flow",
     "The eight stages from terrain loading through to output export."),
    ("06", "System Variables",
     "Inputs, prevention parameters, terrain arrays, and storm presets."),
    ("07", "Assumptions",
     "Modeling simplifications, scope limitations, and what is not simulated."),
    ("08", "Simulation Setup",
     "Software requirements, inputs, run procedures, and numerical config."),
    ("09", "Results",
     "Risk zonation, return-period analysis, and the 24-cell master dataset."),
    ("10", "Conclusion",
     "Findings, observed effectiveness of prevention measures, and outlook."),
]


def build():
    print(f"Building presentation: {OUT_PPT.name}")
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    master_rows = load_master_csv()
    stats       = load_flood_statistics()

    # ------------------------------------------------------------------------
    # First pass: layout the slides; we need total count BEFORE writing them
    # ------------------------------------------------------------------------
    plan: list = []
    plan.append(("cover",))                 # 1
    plan.append(("toc",))                   # 2
    for sec_num, sec_title, blurb in SECTIONS:
        plan.append(("divider", sec_num, sec_title, blurb))

        if sec_num == "01":
            plan.append(("content", sec_num, sec_title, "Project Context",
                         "intro_context"))
            plan.append(("content", sec_num, sec_title, "Three Progressive Builds + Master",
                         "intro_builds"))
        elif sec_num == "02":
            plan.append(("content", sec_num, sec_title, "Why Jade Valley Floods",
                         "problem_why"))
            plan.append(("content", sec_num, sec_title, "Questions This System Answers",
                         "problem_questions"))
        elif sec_num == "03":
            plan.append(("content", sec_num, sec_title, "Six Objectives",
                         "objectives"))
        elif sec_num == "04":
            plan.append(("content", sec_num, sec_title, "Architecture: Five Logical Layers",
                         "arch_layers"))
            plan.append(("content", sec_num, sec_title, "Build Hierarchy",
                         "build_hierarchy"))
            plan.append(("content", sec_num, sec_title, "Core Components",
                         "core_components"))
            plan.append(("content", sec_num, sec_title, "Data Flow",
                         "data_flow"))
        elif sec_num == "05":
            plan.append(("content", sec_num, sec_title, "Stages 1–4: Load, Input, Prevention, Preprocess",
                         "flow_1234"))
            plan.append(("content", sec_num, sec_title, "Stage 5: Per-Timestep Loop",
                         "flow_5"))
            plan.append(("content", sec_num, sec_title, "Stages 6–8: Baseline, Render, Export",
                         "flow_678"))
        elif sec_num == "06":
            plan.append(("content", sec_num, sec_title, "User Input Variables",
                         "vars_user"))
            plan.append(("content", sec_num, sec_title, "Prevention Measure Variables",
                         "vars_prev"))
            plan.append(("content", sec_num, sec_title, "Storm Scenario Presets",
                         "vars_scenarios"))
        elif sec_num == "07":
            plan.append(("content", sec_num, sec_title, "Modeling Assumptions",
                         "assumptions"))
        elif sec_num == "08":
            plan.append(("content", sec_num, sec_title, "Software + Inputs",
                         "setup_software"))
            plan.append(("content", sec_num, sec_title, "Running the System",
                         "setup_running"))
            plan.append(("content", sec_num, sec_title, "Numerical Configuration",
                         "setup_numeric"))
        elif sec_num == "09":
            plan.append(("content", sec_num, sec_title, "HAND Risk Zones",
                         "results_risk"))
            plan.append(("content", sec_num, sec_title, "Return Period Analysis",
                         "results_return"))
            plan.append(("table",    sec_num, sec_title, "Master Results: All 24 Runs", "master"))
            plan.append(("raw_text", sec_num, sec_title, "Master Results Summary (Raw Export)", "master_summary_txt"))
            plan.append(("table",    sec_num, sec_title, "Reduction vs Baseline", "reduction"))
        elif sec_num == "10":
            plan.append(("content", sec_num, sec_title, "Findings",
                         "concl_findings"))
            plan.append(("content", sec_num, sec_title, "Closing Notes",
                         "concl_close"))

    plan.append(("end",))                  # final

    total = len(plan)

    # ------------------------------------------------------------------------
    # Second pass: actually build slides
    # ------------------------------------------------------------------------
    for i, entry in enumerate(plan, start=1):
        kind = entry[0]
        if kind == "cover":
            slide_cover(prs, total)
        elif kind == "toc":
            slide_toc(prs, [(s[0], s[1]) for s in SECTIONS], i, total)
        elif kind == "divider":
            slide_section_divider(prs, entry[1], entry[2], entry[3], i, total)
        elif kind == "content":
            _, sec_num, sec_title, slide_title, key = entry
            paragraphs = CONTENT[key]
            mono_body = key in ("vars_scenarios", "build_hierarchy", "core_components",
                                "data_flow", "setup_numeric")
            slide_text_content(prs, sec_num, sec_title, slide_title,
                               paragraphs, i, total, body_size=13,
                               mono_body=mono_body)
        elif kind == "table":
            _, sec_num, sec_title, slide_title, key = entry
            if key == "master":
                headers = ["Scenario", "Config", "Peak %", "Peak ha",
                           "Max Depth mm", "River %", "Onset min"]
                rows = []
                for r in master_rows:
                    rows.append([
                        _short_scen(r["scenario"]),
                        r["config"],
                        r["peak_flooded_pct"],
                        r["peak_flooded_ha"],
                        int(float(r["peak_depth_mm"])),
                        r["peak_river_pct"],
                        r["overflow_min"] or "—",
                    ])
                slide_table(prs, sec_num, sec_title, slide_title,
                            headers, rows, i, total,
                            col_widths=[Inches(2.7), Inches(1.9),
                                        Inches(1.1), Inches(1.2),
                                        Inches(1.6), Inches(1.1),
                                        Inches(1.4)],
                            header_size=10, body_size=8.5,
                            caption=("Source: Results/data/master_results_table.csv  "
                                     "— Generated by flood_simulation_100%.py "
                                     "(6 PAGASA scenarios × 4 prevention configs)."))
            elif key == "reduction":
                headers = ["Scenario", "Config", "Δ Flooded %",
                           "Δ Flooded ha", "Δ Depth mm", "Reduction %"]
                rows = []
                # Compute deltas inline
                by_scen: dict = {}
                for r in master_rows:
                    by_scen.setdefault(r["scenario"], {})[r["config"]] = r
                for scen, cfgs in by_scen.items():
                    if "Baseline" not in cfgs:
                        continue
                    base = cfgs["Baseline"]
                    bpct = float(base["peak_flooded_pct"])
                    for cfg_label, r in cfgs.items():
                        if cfg_label == "Baseline":
                            continue
                        dpct = float(r["peak_flooded_pct"]) - bpct
                        dha  = float(r["peak_flooded_ha"])  - float(base["peak_flooded_ha"])
                        ddep = float(r["peak_depth_mm"])    - float(base["peak_depth_mm"])
                        red  = (dpct / bpct * 100) if bpct > 0 else 0.0
                        rows.append([
                            _short_scen(scen),
                            cfg_label,
                            f"{dpct:+.2f}",
                            f"{dha:+.2f}",
                            f"{int(ddep):+d}",
                            f"{red:+.1f}%",
                        ])
                slide_table(prs, sec_num, sec_title, slide_title,
                            headers, rows, i, total,
                            col_widths=[Inches(2.9), Inches(2.0),
                                        Inches(1.5), Inches(1.7),
                                        Inches(1.6), Inches(1.4)],
                            header_size=10, body_size=8.5,
                            caption=("Negative values = reduction relative to baseline "
                                     "(good). Light Rain produces no baseline flooding, "
                                     "so reductions are 0."))
        elif kind == "raw_text":
            _, sec_num, sec_title, slide_title, key = entry
            if key == "master_summary_txt":
                slide_raw_text(
                    prs, sec_num, sec_title, slide_title,
                    DATA_DIR / "master_results_summary.txt",
                    i, total,
                    caption=("Verbatim export of Results/data/master_results_summary.txt "
                             "— produced by flood_simulation_100%.py "
                             "(COLLECT ALL DATA / --matrix-only)."),
                    body_size=8.5,
                )
        elif kind == "end":
            slide_end(prs, i, total)

    OUT_PPT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPT))
    print(f"  Saved: {OUT_PPT}")
    print(f"  Total slides: {total}")


def _short_scen(s):
    s = s.replace("Typhoon Signal", "Sig")
    s = s.replace("Tropical Depression", "TD").replace("Tropical Storm", "TS")
    s = s.replace("Severe Typhoon", "Severe")
    s = s.replace(" (", " (")
    return s


def slide_end(prs, slide_num, total):
    slide = make_blank(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=BLACK)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(12.1), Inches(1.2),
                "END OF PRESENTATION", size=48, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(5.6), Inches(3.55), Inches(2.1), Emu(28575),
             fill=WHITE)
    add_textbox(slide, Inches(0.6), Inches(3.7), Inches(12.1), Inches(0.6),
                "JADE VALLEY FLOOD SIMULATION",
                size=18, color=GREY_L, align=PP_ALIGN.CENTER, mono=True)
    add_textbox(slide, Inches(0.6), Inches(4.2), Inches(12.1), Inches(0.5),
                "Davao City, Philippines  |  Final Project",
                size=14, color=GREY_L, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.4),
                "See  Documents/USER_GUIDE.md  and  Documents/User_Guide.docx  "
                "for installation and usage instructions.",
                size=11, color=GREY_M, align=PP_ALIGN.CENTER, mono=True)
    add_textbox(slide, Inches(0.6), Inches(7.10), Inches(8), Inches(0.3),
                "JADE VALLEY FLOOD SIMULATION   |   DAVAO CITY",
                size=8, bold=True, color=GREY_M, mono=True)
    add_textbox(slide, Inches(10.0), Inches(7.10), Inches(2.7), Inches(0.3),
                f"{slide_num} / {total}",
                size=8, bold=True, color=GREY_M, mono=True,
                align=PP_ALIGN.RIGHT)


# =============================================================================
# CONTENT BLOCKS (text material drawn from documentation)
# =============================================================================

CONTENT = {
    "intro_context": [
        "Jade Valley Subdivision is a residential community in Davao City, "
        "Philippines, situated on low-lying terrain that is highly susceptible "
        "to flooding during heavy rainfall events and tropical cyclones.",
        "",
        "The study area covers approximately 326.44 hectares with elevations "
        "ranging from 2.00 m to 71.75 m above mean sea level, mean elevation "
        "of 20.74 m. With 72.3% of the land classified as flat (slope < 5°), "
        "surface runoff accumulates rapidly during storms, exposing the "
        "subdivision to both pluvial flooding from direct rainfall and "
        "fluvial flooding from nearby river overflow.",
        "",
        "This system is an interactive, terrain-aware flood simulation built "
        "entirely in Python (rasterio, NumPy, SciPy, Matplotlib) operating on "
        "a GeoTIFF Digital Elevation Model (DEM) derived from 3D topographic "
        "survey data of the subdivision.",
    ],
    "intro_builds": [
        ("25%",  "Baseline simulation. No prevention measures. Establishes "
                 "the unmitigated flood behaviour as the reference for comparison."),
        ("50%",  "Adds two prevention measures: Riverbank Floodwall and "
                 "Drainage Canal Network. Runs an improved and a baseline "
                 "simulation in parallel and reports the delta."),
        ("75%",  "Adds two further measures: Retention Basin and Elevated "
                 "Emergency Road. Renders all four infrastructure overlays "
                 "on the animated map."),
        ("100%", "Complete system. Same GUI as 75% (storm tab + four "
                 "prevention measures) PLUS a COLLECT ALL DATA button that "
                 "runs the full 6 × 4 = 24 scenario/config matrix and "
                 "exports the master data tables that drive this presentation."),
    ],
    "problem_why": [
        "Davao City is among the most flood-prone urban centres in the "
        "Philippines: proximity to river systems, high annual rainfall, and "
        "rapid urbanisation of low-lying areas.",
        "",
        "Jade Valley exemplifies the compounding flood risk in such "
        "communities: flat terrain limits natural drainage while proximity "
        "to river channels exposes residents to both surface runoff "
        "accumulation and fluvial overflow during storm events.",
        "",
        "Despite these conditions, flood risk planning in residential "
        "subdivisions often lacks site-specific quantitative modeling — "
        "meaning emergency preparedness, infrastructure investments, and "
        "land use recommendations are made without direct evidence of how "
        "water actually moves across the terrain during real storms.",
    ],
    "problem_questions": [
        ("Q1", "Which specific zones flood first, and at what rainfall "
               "threshold does inundation become critical?"),
        ("Q2", "How far does river overflow spread across the subdivision "
               "during typhoons of varying intensities?"),
        ("Q3", "How much flood extent and peak depth reduction can be "
               "achieved by constructing prevention infrastructure — and "
               "which combination provides the best protection?"),
        ("Q4", "Which areas can safely serve as evacuation zones during a "
               "100-year flood event?"),
        "",
        "This simulation system answers each of these questions with a "
        "terrain-aware, scenario-driven model grounded in standard "
        "hydrological methods.",
    ],
    "objectives": [
        ("3.1", "Terrain-Accurate Flood Inundation Modeling — D8 flow routing, "
                "Green-Ampt infiltration, BFS overflow dilation."),
        ("3.2", "Multi-Scenario Storm Coverage — PAGASA-calibrated scenarios "
                "from Light Rain (15 mm / 2 h) to Typhoon Signal 3 "
                "(400 mm / 18 h)."),
        ("3.3", "Flood Risk Quantification — five-class HAND-based risk "
                "zonation (Very High / High / Medium / Low / Safe)."),
        ("3.4", "Prevention Measure Evaluation — frame-by-frame comparison "
                "of baseline vs. improved runs across four mitigation "
                "configurations."),
        ("3.5", "Interactive Visualisation — animated multi-layer figure "
                "with playback controls."),
        ("3.6", "Data Export — GIF animations, time-series CSV files, "
                ".npy terrain arrays, JSON statistics, and human-readable "
                "reports."),
    ],
    "arch_layers": [
        ("LAYER 1", "Data Layer  |  Reads the GeoTIFF DEM and JPEG satellite "
                    "background from disk. All terrain computations derive "
                    "from the 57 × 61 elevation grid."),
        ("LAYER 2", "Input Layer  |  A Tkinter GUI collects storm parameters "
                    "and prevention settings (the 100% build skips this and "
                    "runs the matrix headlessly)."),
        ("LAYER 3", "Preprocessing Layer  |  Applies prevention DEM "
                    "modifications, then runs the hydrological pipeline: "
                    "depression fill, D8 flow direction, flow accumulation, "
                    "slope, runoff coefficient, infiltration capacity, bank "
                    "elevation map."),
        ("LAYER 4", "Simulation Engine  |  Frame-by-frame physics loop: "
                    "rainfall addition, water routing, river overflow via "
                    "BFS dilation, infiltration, drainage. Parallel baseline "
                    "sim runs simultaneously when prevention is active."),
        ("LAYER 5", "Output Layer  |  Matplotlib renders the animation; "
                    "results export as GIF animations and CSV/JSON/.npy "
                    "artefacts."),
    ],
    "build_hierarchy": [
        "BUILD          SCRIPT                          PREVENTION MEASURES",
        "-" * 80,
        "25%  Baseline  flood_simulation_25%.py          None",
        "50%  Two-meas. flood_simulation_50%.py          Floodwall + Drainage Canal",
        "75%  Four-meas flood_simulation_75%.py          Floodwall + Canal +",
        "                                                Retention Basin + Elevated Road",
        "100% Master    flood_simulation_100%.py         Headless 6 × 4 matrix runner",
        "                                                + HAND/risk-zone analysis",
    ],
    "core_components": [
        "load_dem()                  Reads GeoTIFF; fills NoData; returns (dem, cellsize)",
        "_fill_depressions()         Priority-Flood sink removal",
        "_d8_flow_direction()        Steepest-descent D8 routing",
        "_flow_accumulation()        Topological-sort upstream count",
        "build_stream_mask()         Top-8% accumulation = river channel",
        "apply_floodwall()           DEM modification: raises western bank cells",
        "apply_drainage_canal()      DEM modification: lowers east/south corridors",
        "apply_retention_basin()     DEM modification: excavates basin depression [75%]",
        "apply_elevated_road()       DEM modification: raises emergency road berm [75%]",
        "apply_prevention_measures() Combines active measures into one modified DEM",
        "FloodSimulation             Class — holds state; runs physics per timestep",
        "    .add_rainfall()         Runoff-weighted rainfall onto grid",
        "    .route_water()          D8 surface-water diffusion (12 sub-iters)",
        "    .apply_river_overflow() BFS dilation of flood front from channel",
        "    .apply_infiltration()   Green-Ampt inspired soil absorption",
        "    .apply_drainage()       Channel + canal drainage removal",
        "    .step()                 Calls all five methods for one timestep",
        "wind_rainfall_map()         Wind-modified spatial rainfall weights",
        "run_simulation()            Orchestrates preprocessing, loop, rendering",
        "SimulationGUI               Tkinter two-tab GUI",
    ],
    "data_flow": [
        "GeoTIFF DEM",
        "    |",
        "    +-- Prevention DEM modifications (if enabled)",
        "    |",
        "    +-- Hydrological preprocessing  -->  fdir, accum, river_mask,",
        "    |                                    slope, runoff_coeff,",
        "    |                                    max_inf, bank_elev",
        "    |",
        "    +-- [Improved Sim]  FloodSimulation(modified_dem)  -+",
        "    |                   frame loop x num_frames         +-- Animation",
        "    +-- [Baseline Sim]  FloodSimulation(original_dem)  -+   + CSV Export",
    ],
    "flow_1234": [
        ("STAGE 1", "Data Loading  |  Read GeoTIFF DEM (rasterio) and JPEG "
                    "background (Pillow). Resize JPEG to ≤ 2048 px on longest "
                    "side. Cell size corrected via cosine-latitude scaling."),
        ("STAGE 2", "User Input  |  Two-tab Tkinter GUI gathers storm "
                    "parameters and prevention settings. 100% build skips "
                    "this stage and uses CLI flags / defaults."),
        ("STAGE 3", "DEM Modification  |  For each active measure, modify "
                    "elevations on target cells: raise floodwall, lower "
                    "canals, excavate retention basin, raise emergency road."),
        ("STAGE 4", "Hydrological Preprocessing  |  Priority-Flood fill, "
                    "D8 flow direction, flow accumulation, stream mask, "
                    "slope, runoff coefficient, infiltration capacity, "
                    "bank-elevation map."),
    ],
    "flow_5": [
        "The simulation loop runs num_frames = ceil(duration_h / timestep_h) "
        "iterations. At each timestep, six physics operations execute in order:",
        "",
        ("STEP 1", "Rainfall intensity factor scales rate by storm pattern "
                   "(uniform / progressive / burst / decreasing)."),
        ("STEP 2", "add_rainfall(): rain × intensity × dt × runoff_coeff "
                   "× wind_map."),
        ("STEP 3", "route_water(): D8 surface diffusion, 12 sub-iterations."),
        ("STEP 4", "apply_river_overflow(): rises river level; BFS dilation "
                   "expands flood front; floodwall blocks overflow until "
                   "overtopped; canal cells receive 40% less river water."),
        ("STEP 5", "apply_infiltration(): saturation grows; infiltration "
                   "removes water (Green-Ampt inspired)."),
        ("STEP 6", "apply_drainage(): base + channel + canal drainage; "
                   "shallow water decays."),
    ],
    "flow_678": [
        ("STAGE 6", "Parallel Baseline Run  |  When prevention is enabled, "
                    "an identical sim runs on the unmodified DEM with no "
                    "measures; deltas feed the live stats panel."),
        ("STAGE 7", "Animation Rendering  |  Five composited Matplotlib "
                    "layers: satellite background, river channel band, "
                    "rain depth, river overflow, prevention overlays."),
        ("STAGE 8", "Output Export  |  GIF animation, time-series CSVs, "
                    "and (via the 100% build) master_results_table.csv/.json, "
                    "master_results_summary.txt, master_reduction_table.txt, "
                    "plus the HAND/risk-zone artefacts."),
    ],
    "vars_user": [
        ("rainfall_mm",   "0–600 mm     |  Total storm rainfall."),
        ("duration_h",    "0.5–24.0 h   |  Storm duration."),
        ("timestep_min",  "1–60 min     |  Simulation timestep."),
        ("wind_speed",    "0–200 km/h   |  Wind speed."),
        ("wind_dir",      "0–360°       |  Wind direction."),
        ("soil_sat_pct",  "0–100 %      |  Initial soil saturation."),
        ("drain_cap",     "0–30 mm/hr   |  Baseline drainage capacity."),
        ("pattern",       "uniform / progressive / burst / decreasing"),
        ("start_time",    "HH:MM        |  Clock time at frame 0."),
    ],
    "vars_prev": [
        ("use_floodwall / wall_height_m",
         "Riverbank Floodwall, crest 0.5–5.0 m (default 1.5 m)."),
        ("use_canal / canal_depth_m",
         "Drainage Canal Network, depth 0.5–5.0 m (default 2.0 m)."),
        ("use_basin / basin_depth_m   [75% / 100%]",
         "Retention Basin, depth 2.0–10.0 m (default 6.0 m)."),
        ("use_road / road_height_m   [75% / 100%]",
         "Elevated Emergency Road, raise 0.5–3.0 m (default 1.5 m)."),
    ],
    "vars_scenarios": [
        "ID  NAME                                  RAINFALL   DURATION  PATTERN",
        "-" * 80,
        " 1  Light Rain                              15 mm       2.0 h  uniform",
        " 2  Moderate Rain                           36 mm       3.0 h  progressive",
        " 3  Heavy Rain                              90 mm       4.0 h  progressive",
        " 4  Typhoon Signal 1 (Tropical Depression) 150 mm       8.0 h  progressive",
        " 5  Typhoon Signal 2 (Tropical Storm)      250 mm      12.0 h  burst",
        " 6  Typhoon Signal 3 (Severe Typhoon)      400 mm      18.0 h  burst",
        " 7  Custom                                 user def   user def  user def",
    ],
    "assumptions": [
        ("5.1", "Single-Direction Flow — D8: each cell routes to one of "
                "8 neighbours (steepest descent)."),
        ("5.2", "Static DEM — no erosion/sedimentation; prevention measures "
                "are applied once before t=0."),
        ("5.3", "Uniform Initial Soil Saturation — spatially uniform; "
                "infiltration capacity uses elevation as proxy."),
        ("5.4", "Simplified Green-Ampt Infiltration — scalar approximation; "
                "no explicit wetting front."),
        ("5.5", "BFS River Overflow — efficient dilation; not a full "
                "shallow-water equation solver."),
        ("5.6", "Instantaneous Prevention Deployment — measures are fully "
                "effective from t=0."),
        ("5.7", "No Groundwater Interaction — surface hydrology only."),
        ("5.8", "Wind Uniformity — constant throughout the storm."),
        ("5.9", "PAGASA-Calibrated Scenarios — Mindanao / Davao classification."),
        ("5.10", "Retention Basin — instantaneous fill, no inlet/spillway model."),
        ("5.11", "Grid Resolution — ~30.64 m / cell; sub-grid features "
                 "approximated via runoff coefficient and drainage capacity."),
    ],
    "setup_software": [
        ("Python",        "3.10 or higher"),
        ("Virtual env",   ".venv (project root)"),
        ("Required pkgs", "numpy  scipy  matplotlib  rasterio  Pillow  "
                          "ezdxf  pysheds  pandas  python-pptx  python-docx"),
        "",
        ("Required",      "Map Topography/3D/JVS_Simulation.tif    (GeoTIFF DEM)"),
        ("Required",      "Map Topography/2D/JVS_2D.jpg            (background)"),
        ("Optional",      "Map Topography/2D/...2D_vectorial.dxf   (vector overlay)"),
        ("Optional",      "Map Topography/3D/...3D_modeling.dxf    (3D reference)"),
    ],
    "setup_running": [
        ("25% (GUI)",        'python "Main/flood_simulation_25%.py"             # baseline only'),
        ("50% (GUI)",        'python "Main/flood_simulation_50%.py"             # + Wall + Canal'),
        ("75% (GUI)",        'python "Main/flood_simulation_75%.py"             # + Basin + Road'),
        ("100% (GUI+matrix)",'python "Main/flood_simulation_100%.py"            # full system'),
        ("100% headless",    'python "Main/flood_simulation_100%.py" --matrix-only'),
        ("HAND analysis",    'python "Main/hand_risk_analysis.py"               # idempotent'),
        ("Reality check",    'python "Main/run_validation.py"                   # vs historical events'),
    ],
    "setup_numeric": [
        "DEM Grid                  : 57 rows × 61 cols",
        "Cell Size                 : 30.64 m",
        "Total Study Area          : 326.44 ha",
        "Elevation Range           : 2.00 m – 71.75 m",
        "Mean Elevation            : 20.74 m",
        "Mean Slope                : 4.3°",
        "Stream Threshold          : Top 8% of flow accumulation (92nd pctile)",
        "HAND Smooth Sigma         : 0.5",
        "Curve Number              : 85  (urban mix reference)",
        "Routing Sub-steps         : 12 D8 iterations per timestep",
        "River Overflow BFS Hops   : 1 (light) to 4-6 (typhoon)",
        "Matrix Timestep           : 20 min  (100% build default)",
    ],
    "results_risk": [
        "The HAND (Height Above Nearest Drainage) model classifies every "
        "terrain cell into one of five risk zones:",
        "",
        ("VERY HIGH",  "26.66 ha (8.2%)   — Floods in Heavy Rain ≥ 90 mm. "
                       "Mandatory evacuation during typhoons."),
        ("HIGH",       "31.83 ha (9.7%)   — Floods in Typhoon Signal 1 "
                       "≥ 150 mm. Pre-evacuation alert."),
        ("MEDIUM",     "32.30 ha (9.9%)   — Floods in Typhoon Signal 2 "
                       "≥ 250 mm. Monitor advisories."),
        ("LOW",        "35.86 ha (11.0%)  — Floods only in extreme events "
                       "(Sig 3 / ≥ 400 mm). Standby."),
        ("SAFE",       "174.81 ha (53.6%) — Above 100-yr flood level. "
                       "Suitable as evacuation destination."),
        "",
        "Reference: Nobre et al. (2011), Journal of Hydrology 404:13–29.",
    ],
    "results_return": [
        ("5-YEAR",    "HAND threshold 1.0 m  —  51.64 ha (15.8%) flooded  —  "
                      "mean depth 0.73 m."),
        ("10-YEAR",   "HAND threshold 2.0 m  —  70.04 ha (21.5%) flooded  —  "
                      "mean depth 1.41 m."),
        ("25-YEAR",   "HAND threshold 3.5 m  —  90.03 ha (27.6%) flooded  —  "
                      "mean depth 2.43 m."),
        ("100-YEAR",  "HAND threshold 6.0 m  —  115.76 ha (35.5%) flooded  —  "
                      "mean depth 4.12 m."),
        "",
        "These results underline the severity of the 100-year design event "
        "(35.5% of the subdivision inundated) and the rationale for the "
        "layered prevention strategy evaluated by the simulation matrix.",
    ],
    "concl_findings": [
        "The 100% master simulation matrix demonstrates a consistent, "
        "monotonic improvement as additional prevention measures are layered:",
        "",
        ("OBSERVATION 1",
         "Light Rain (15 mm / 2 h) produces no measurable flooding under "
         "any configuration — the subdivision can absorb light convective "
         "events without intervention."),
        ("OBSERVATION 2",
         "Wall+Canal alone reduces peak flooded area by 5–24% across the "
         "moderate-to-heavy band. The retention basin and elevated road "
         "extend the reduction to 11–41% (full prevention) and 13–41% "
         "(large-scale prevention) — the basin contributes the largest "
         "single reduction at moderate intensities."),
        ("OBSERVATION 3",
         "Peak depth can rise slightly with prevention: water concentrates "
         "in the basin sink while spread shrinks. The relevant decision "
         "metric is therefore extent (% / ha), not depth alone."),
        ("OBSERVATION 4",
         "Onset time is unchanged by prevention measures — they reduce "
         "extent, not timing. Early-warning systems remain essential."),
    ],
    "concl_close": [
        "The Jade Valley Flood Simulation system successfully models the "
        "hydrological behaviour of a real-world Philippine subdivision "
        "under the full spectrum of PAGASA-classified storm conditions.",
        "",
        "By combining D8 terrain-based flow routing, scenario-adaptive "
        "river overflow via BFS dilation, and Green-Ampt-inspired "
        "infiltration, the model captures the essential dynamics of urban "
        "flooding at a computationally accessible scale.",
        "",
        "The three-build hierarchy plus the 100% master orchestrator "
        "provides quantitative, frame-by-frame comparison of flood extent "
        "with and without each prevention measure. The master CSV/JSON "
        "exports support direct use in infrastructure planning, "
        "community disaster preparedness, and policy recommendations for "
        "flood-vulnerable residential areas such as Jade Valley.",
    ],
}


if __name__ == "__main__":
    build()
